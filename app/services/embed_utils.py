"""
Document processing utilities for PDF and PPTX files.
Handles text extraction, chunking, embedding generation, and ChromaDB storage.
Uses OpenAI text-embedding-3-small model for high-quality embeddings.
"""

import chromadb
from chromadb.config import Settings as ChromaSettings
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document
import PyPDF2
from pptx import Presentation
from typing import List, Tuple
import os

from app.config import settings


def extract_text_from_pptx(file_path: str) -> List[Tuple[int, str]]:
    """
    Extract text from a PowerPoint (PPTX) file, preserving slide numbers.

    Args:
        file_path: Path to the PPTX file

    Returns:
        List of (slide_number, slide_text) tuples
    """
    prs = Presentation(file_path)
    slides = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_parts.append(shape.text)
        slide_text = "\n".join(text_parts)
        if slide_text.strip():
            slides.append((slide_num, slide_text))

    return slides


def extract_text_from_pdf(file_path: str) -> List[Tuple[int, str]]:
    """
    Extract text from a PDF file, preserving page numbers.

    Args:
        file_path: Path to the PDF file

    Returns:
        List of (page_number, page_text) tuples
    """
    pages = []

    with open(file_path, 'rb') as file:
        pdf_reader = PyPDF2.PdfReader(file)

        for page_num, page in enumerate(pdf_reader.pages, start=1):
            text = page.extract_text()
            if text and text.strip():
                pages.append((page_num, text))

    return pages


def chunk_text(text: str) -> List[Document]:
    """
    Split text into smaller chunks using RecursiveCharacterTextSplitter.
    (Legacy function - kept for backward compatibility)

    Args:
        text: Input text to split

    Returns:
        List of Document objects containing text chunks
    """
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len,
    )

    docs = splitter.create_documents([text])
    return docs


def chunk_by_page_or_slide(
    pages_or_slides: List[Tuple[int, str]],
    document_type: str,
    source_file: str,
    collection_name: str,
    chunk_size: int = 1500,
    chunk_overlap: int = 150
) -> List[Document]:
    """
    Create chunks from a file by combining all pages/slides and chunking with overlap.
    Each file represents one lecture - extract lecture number from filename and assign
    it as slide_number to all chunks.

    Args:
        pages_or_slides: List of (page/slide_number, text) tuples
        document_type: "pdf" or "pptx"
        source_file: Original filename (e.g., "Lecture 1_slides.pdf")
        collection_name: ChromaDB collection name
        chunk_size: Max characters per chunk (default: 1500)
        chunk_overlap: Character overlap between chunks (default: 150)

    Returns:
        List of Document objects with slide_number metadata (lecture number)
    """
    import re
    
    # Extract lecture number from filename (e.g., "Lecture 1_slides.pdf" -> 1)
    lecture_match = re.search(r'Lecture\s+(\d+)', source_file, re.IGNORECASE)
    if lecture_match:
        slide_number = int(lecture_match.group(1))
    else:
        # Fallback: use 0 if can't extract lecture number
        slide_number = 0
    
    # Combine all pages/slides into one continuous text string
    all_text = "\n\n".join([text for _, text in pages_or_slides])
    
    # Chunk the combined text with RecursiveCharacterTextSplitter
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        length_function=len,
    )
    
    chunks = splitter.create_documents([all_text])
    
    # Create Document objects with slide_number metadata
    all_docs = []
    for chunk_idx, chunk in enumerate(chunks):
        metadata = {
            "source_file": source_file,
            "document_type": document_type,
            "slide_number": slide_number,  # Lecture number from filename
            "chunk_index": chunk_idx,
            "collection_name": collection_name
        }
        
        doc = Document(
            page_content=chunk.page_content,
            metadata=metadata
        )
        all_docs.append(doc)
    
    return all_docs


def store_in_chroma(
    docs: List[Document],
    collection_name: str = "default_collection"
):
    """
    Store document chunks in ChromaDB with embeddings.
    Uses OpenAI text-embedding-3-small model for high-quality embeddings.

    Args:
        docs: List of Document objects to store
        collection_name: Name of the ChromaDB collection

    Returns:
        ChromaDB collection instance
        
    Raises:
        ValueError: If OpenAI API key is not configured
    """
    
    # Ensure ChromaDB directory exists
    os.makedirs(settings.chroma_db_path, exist_ok=True)

    # Check for OpenAI API key
    if not settings.openai_api_key:
        raise ValueError(
            "OpenAI API key is required for embeddings. "
            "Please set OPENAI_API_KEY in your .env file."
        )

    # Use singleton ChromaDB client to prevent conflicts
    # Import ChromaService to reuse its singleton client
    from app.services.chroma_service import ChromaService
    
    try:
        # Get singleton client instance
        chroma_service = ChromaService()
        client = chroma_service.client
        
        # Use OpenAI embedding function
        from chromadb.utils.embedding_functions import OpenAIEmbeddingFunction
        embedding_function = OpenAIEmbeddingFunction(
            api_key=settings.openai_api_key,
            model_name="text-embedding-3-small"
        )

        collection = client.get_or_create_collection(
            name=collection_name,
            embedding_function=embedding_function
        )
    except Exception as e:
        import logging
        logger = logging.getLogger(__name__)
        logger.error(f"Error creating ChromaDB client/collection: {e}")
        raise

    # Prepare data for ChromaDB
    documents = [doc.page_content for doc in docs]
    
    # Generate unique IDs that include source_file to prevent overwrites
    # Format: {collection_name}_{source_file_hash}_{chunk_index}
    import hashlib
    ids = []
    seen_ids = set()  # Track IDs to ensure uniqueness
    
    for i, doc in enumerate(docs):
        source_file = doc.metadata.get("source_file", "unknown") if doc.metadata else "unknown"
        # Create a short hash of source_file to keep IDs manageable
        source_hash = hashlib.md5(source_file.encode()).hexdigest()[:8]
        chunk_id = f"{collection_name}_{source_hash}_{i}"
        
        # Ensure ID is unique (add suffix if collision detected)
        original_id = chunk_id
        counter = 0
        while chunk_id in seen_ids:
            counter += 1
            chunk_id = f"{original_id}_{counter}"
        
        seen_ids.add(chunk_id)
        ids.append(chunk_id)

    # Ensure each document has non-empty metadata (ChromaDB requirement)
    metadatas = []
    for i, doc in enumerate(docs):
        if doc.metadata and len(doc.metadata) > 0:
            metadatas.append(doc.metadata)
        else:
            # Provide default metadata if empty
            metadatas.append({"chunk_index": i, "source": collection_name})

    # Add documents (embeddings generated automatically by ChromaDB)
    import logging
    logger = logging.getLogger(__name__)
    
    # Validate data before adding
    if len(documents) != len(ids) or len(documents) != len(metadatas):
        error_msg = f"Data length mismatch: {len(documents)} documents, {len(ids)} ids, {len(metadatas)} metadatas"
        logger.error(f"   ✗ {error_msg}")
        raise ValueError(error_msg)
    
    # Check for duplicate IDs
    if len(ids) != len(set(ids)):
        duplicates = [id for id in ids if ids.count(id) > 1]
        error_msg = f"Duplicate IDs detected: {set(duplicates)}"
        logger.error(f"   ✗ {error_msg}")
        raise ValueError(error_msg)
    
    logger.info(f"   Adding {len(documents)} documents to ChromaDB collection '{collection_name}'...")
    logger.info(f"   Generating embeddings using OpenAI text-embedding-3-small...")
    logger.debug(f"   Sample IDs (first 3): {ids[:3] if len(ids) >= 3 else ids}")
    
    try:
        collection.add(
            documents=documents,
            ids=ids,
            metadatas=metadatas
        )
        logger.info(f"   ✓ Successfully added {len(documents)} documents with embeddings")
        logger.debug(f"   Collection '{collection_name}' now has {collection.count()} total documents")
    except Exception as e:
        logger.error(f"   ✗ Error adding documents to ChromaDB: {str(e)}")
        logger.error(f"   Collection: {collection_name}, Documents count: {len(documents)}")
        logger.error(f"   First few IDs: {ids[:5] if len(ids) >= 5 else ids}")
        raise

    return collection


def process_document(
    file_path: str,
    file_type: str,
    collection_name: str = "default_collection"
) -> dict:
    """
    Complete pipeline to process a document: extract text, chunk, embed, and store.
    Uses one chunk per page/slide approach for citation support.

    Args:
        file_path: Path to the document file
        file_type: Type of file ('pdf' or 'pptx')
        collection_name: Name for the ChromaDB collection

    Returns:
        Dictionary with processing results
    """
    source_file = os.path.basename(file_path)
    
    # Extract text with page/slide tracking
    if file_type == "pdf":
        pages = extract_text_from_pdf(file_path)
        if not pages:
            raise ValueError("No text could be extracted from the document")
        docs = chunk_by_page_or_slide(
            pages, "pdf", source_file, collection_name
        )
        total_text_length = sum(len(text) for _, text in pages)
    elif file_type == "pptx":
        slides = extract_text_from_pptx(file_path)
        if not slides:
            raise ValueError("No text could be extracted from the document")
        docs = chunk_by_page_or_slide(
            slides, "pptx", source_file, collection_name
        )
        total_text_length = sum(len(text) for _, text in slides)
    else:
        raise ValueError(f"Unsupported file type: {file_type}")

    # Store in ChromaDB (this will also generate embeddings)
    vector_store = store_in_chroma(docs, collection_name)

    return {
        "chunks_count": len(docs),
        "collection_name": collection_name,
        "text_length": total_text_length,
    }
